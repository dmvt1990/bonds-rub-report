"""
Слайд «Рублевые облигации — ОФЗ с фиксированным купоном».

Это упрощённая версия стр. 2 образца: без scatter-диаграммы, графика
индексов и макро-блока. Только заголовок и таблица.
"""
from typing import List, Optional, Set

from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm, Pt

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY, COLOR_TEXT_SECONDARY, COLOR_TEXT_MUTED,
    TABLE_SLIDE_TITLE_FONT_SIZE,
    TABLE_SLIDE_LEFT, TABLE_SLIDE_TOP,
    TABLE_SLIDE_SOURCE_TOP, TABLE_SLIDE_SOURCE_FONT_SIZE,
    TABLE_SLIDE_FOOTNOTE_TOP, TABLE_SLIDE_FOOTNOTE_FONT_SIZE,
    TABLE_OFZ_COL_WIDTHS_CM,
)
from ..formatting import (
    format_date_dmy, format_percent, format_number,
    format_integer, format_duration,
)
from ..models import Bond
from ..tables import TableColumn, render_bond_table
from ..fonts import set_run_font


# Используем "Базовый слайд с заголовком, номером страницы и годом"
# (Layout 3) — он у Газпромбанка для содержательных слайдов.
_TABLE_SLIDE_LAYOUT_INDEX = 3

# Заголовок слайда — как в образце ("Рублевые облигации")
SLIDE_TITLE = "Рублевые облигации"

# Заголовок секции таблицы. По запросу: короткое "ОФЗ" + пояснение о типе купона.
SECTION_TITLE_OFZ = "ОФЗ с фиксированным купоном"

# Столбцы таблицы (сокращённый набор — убраны
# "Дата Put/Call", "Дох-ть к оферте", "Мин. лот ₽").
# Все колонки выровнены по центру — так аккуратнее читается компактная
# таблица с разнотипными данными (даты, проценты, числа).
OFZ_COLUMNS = [
    TableColumn("№", TABLE_OFZ_COL_WIDTHS_CM[0], align="center"),
    TableColumn("Выпуск", TABLE_OFZ_COL_WIDTHS_CM[1], align="center"),
    TableColumn("ISIN", TABLE_OFZ_COL_WIDTHS_CM[2], align="center"),
    TableColumn("Дата\nпогашения", TABLE_OFZ_COL_WIDTHS_CM[3], align="center"),
    TableColumn("Купон,\n% год.", TABLE_OFZ_COL_WIDTHS_CM[4], align="center"),
    TableColumn("Цена, %\nот ном.", TABLE_OFZ_COL_WIDTHS_CM[5], align="center"),
    TableColumn("Дох-ть к\nпогаш.,\n% год.", TABLE_OFZ_COL_WIDTHS_CM[6], align="center"),
    TableColumn("Мод.\nдюр.", TABLE_OFZ_COL_WIDTHS_CM[7], align="center"),
    TableColumn("Объём\nвыпуска,\nмлн ₽", TABLE_OFZ_COL_WIDTHS_CM[8], align="center"),
    TableColumn("Период.\nкупона", TABLE_OFZ_COL_WIDTHS_CM[9], align="center"),
    TableColumn("В\nперечне", TABLE_OFZ_COL_WIDTHS_CM[10], align="center"),
]


def _bond_to_row(i: int, b: Bond, highlighted: set) -> List[str]:
    """Одна облигация → список строковых значений для таблицы."""
    mark = "✓" if b.isin in highlighted else "—"
    return [
        str(i),
        b.shortname or "—",
        b.isin,
        format_date_dmy(b.matdate),
        format_percent(b.couponpercent, decimals=2),
        format_number(b.close, decimals=2),
        format_percent(b.yield_close, decimals=2),
        format_duration(b.duration_years),
        format_integer(b.issuesize_mln_rub),
        str(b.couponfrequency) if b.couponfrequency else "—",
        mark,
    ]


def _fill_title_placeholder(slide, text: str) -> bool:
    """
    Заполняет placeholder заголовка на слайде (если он есть).

    Это важно для PowerPoint: пока placeholder заголовка пустой, PowerPoint
    рисует на его месте prompt-текст из layout ("Образец. Название
    слайда..."). Достаточно положить в placeholder любой текст — подсказка
    исчезнет.

    Возвращает True, если placeholder найден и заполнен.
    """
    for ph in slide.placeholders:
        # type "title" (idx=0) — это основной заголовок слайда.
        # В template.pptx у layout 3 он как раз первый.
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            set_run_font(run, FONT_PRIMARY)
            run.font.size = TABLE_SLIDE_TITLE_FONT_SIZE
            run.font.color.rgb = COLOR_TEXT_PRIMARY
            return True
    return False


def _blank_other_placeholders(slide) -> None:
    """
    Опустошает все прочие placeholder'ы (не title), чтобы PowerPoint
    не рисовал на их месте prompt-тексты вроде "Образец. Текст слайда".
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue  # title обрабатывается отдельно
        try:
            tf = ph.text_frame
            tf.clear()
            # Добавим пустой run, чтобы placeholder считался "заполненным"
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ""
        except Exception:
            pass


def _add_source_line(slide, report_date) -> None:
    """Подпись "Источник: Cbonds, 16.04.2026" под таблицей, справа."""
    box = slide.shapes.add_textbox(
        Cm(22.00), TABLE_SLIDE_SOURCE_TOP, Cm(10.70), Cm(0.50),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Источник: Cbonds, {format_date_dmy(report_date)}"
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_SOURCE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_SECONDARY


def _add_footnote(slide) -> None:
    """Мелкая сноска снизу о том, что данные индикативные."""
    box = slide.shapes.add_textbox(
        Cm(1.13), TABLE_SLIDE_FOOTNOTE_TOP, Cm(31.60), Cm(0.40),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Внимание! Приведенные котировки и расчеты являются индикативными "
        "и подлежат регулярному обновлению. Отдельные параметры могут "
        "существенно изменяться под влиянием рыночной конъюнктуры."
    )
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_FOOTNOTE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_MUTED


def render_ofz_fixed(
    prs: PresentationType,
    bonds: List[Bond],
    section_title: str = SECTION_TITLE_OFZ,
    report_date=None,
    highlighted: Optional[set] = None,
) -> None:
    """
    Добавляет слайд с таблицей ОФЗ в презентацию.

    Облигации сортируются по дате погашения (от ближней к дальней),
    независимо от порядка в bonds.yaml — это делает таблицу читаемой
    для инвестора: сначала короткие, потом длинные выпуски.

    Args:
        prs: презентация на основе template.pptx
        bonds: список облигаций, порядок не важен
        section_title: заголовок секции над шапкой колонок.
        report_date: дата отчёта для подписи "Источник: Cbonds, ..."
        highlighted: множество ISIN, которые отмечены галочкой в колонке
                     "В перечне". Если не передано — никто не отмечен.
    """
    highlighted = highlighted or set()

    layout = prs.slide_layouts[_TABLE_SLIDE_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    _fill_title_placeholder(slide, SLIDE_TITLE)
    _blank_other_placeholders(slide)

    from datetime import date as _date
    sorted_bonds = sorted(
        bonds,
        key=lambda b: (b.matdate or _date.max),
    )

    rows = [_bond_to_row(i + 1, b, highlighted) for i, b in enumerate(sorted_bonds)]
    render_bond_table(
        slide,
        left=TABLE_SLIDE_LEFT,
        top=TABLE_SLIDE_TOP,
        columns=OFZ_COLUMNS,
        rows=rows,
        section_title=section_title,
    )

    if report_date is not None:
        _add_source_line(slide, report_date)
    _add_footnote(slide)
