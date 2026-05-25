"""
Слайд «Рублевые облигации — облигации с плавающим купоном».

Отличия от ОФЗ:
- купон задаётся ФОРМУЛОЙ из bonds.yaml ("КС ЦБ РФ +1,5%"), а не числом
- нет колонок YTM и дюрации (MOEX их для флоатеров не считает корректно)
- есть колонка "Дата оферты" — у флоатеров часто есть опцион Put/Call

Если бумаг больше, чем помещается на один слайд (~28 строк), нужно будет
добавить пагинацию — сейчас один слайд.
"""
from typing import Dict, List, Optional

from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm, Pt

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY, COLOR_TEXT_SECONDARY, COLOR_TEXT_MUTED,
    TABLE_SLIDE_TITLE_FONT_SIZE,
    TABLE_SLIDE_LEFT, TABLE_SLIDE_TOP,
    TABLE_SLIDE_SOURCE_TOP, TABLE_SLIDE_SOURCE_FONT_SIZE,
    TABLE_SLIDE_FOOTNOTE_TOP, TABLE_SLIDE_FOOTNOTE_FONT_SIZE,
    TABLE_FLOATERS_COL_WIDTHS_CM,
)
from ..formatting import (
    format_date_dmy, format_number, format_integer,
)
from ..models import Bond
from ..tables import TableColumn, render_bond_table
from ..fonts import set_run_font


# Используем тот же layout 3 (с заголовком, номером страницы и годом)
_TABLE_SLIDE_LAYOUT_INDEX = 3

SLIDE_TITLE = "Рублевые облигации"
SECTION_TITLE_FLOATERS = "Облигации с плавающим купоном"

# 9 колонок (без YTM, дюрации и минимального лота)
FLOATER_COLUMNS = [
    TableColumn("№", TABLE_FLOATERS_COL_WIDTHS_CM[0], align="center"),
    TableColumn("Выпуск", TABLE_FLOATERS_COL_WIDTHS_CM[1], align="center"),
    TableColumn("ISIN", TABLE_FLOATERS_COL_WIDTHS_CM[2], align="center"),
    TableColumn("Дата\nпогашения", TABLE_FLOATERS_COL_WIDTHS_CM[3], align="center"),
    TableColumn("Дата\nоферты", TABLE_FLOATERS_COL_WIDTHS_CM[4], align="center"),
    TableColumn("Купон", TABLE_FLOATERS_COL_WIDTHS_CM[5], align="center"),
    TableColumn("Цена, %\nот ном.", TABLE_FLOATERS_COL_WIDTHS_CM[6], align="center"),
    TableColumn("Объём\nвыпуска,\nмлн ₽", TABLE_FLOATERS_COL_WIDTHS_CM[7], align="center"),
    TableColumn("Период.\nкупона", TABLE_FLOATERS_COL_WIDTHS_CM[8], align="center"),
    TableColumn("В\nперечне", TABLE_FLOATERS_COL_WIDTHS_CM[9], align="center"),
]


def _bond_to_row(i: int, b: Bond, coupon_formula: Optional[str], offer_override: Optional[str], highlighted: set) -> List[str]:
    """
    Одна облигация-флоатер → строка для таблицы.
    """
    if offer_override:
        offer_str = offer_override
    else:
        offer_str = format_date_dmy(b.offerdate)

    mark = "✓" if b.isin in highlighted else "—"

    return [
        str(i),
        b.shortname or "—",
        b.isin,
        format_date_dmy(b.matdate),
        offer_str,
        coupon_formula or "—",
        format_number(b.close, decimals=2),
        format_integer(b.issuesize_mln_rub),
        str(b.couponfrequency) if b.couponfrequency else "—",
        mark,
    ]


def _fill_title_placeholder(slide, text: str) -> bool:
    """Заполняет title placeholder — убирает placeholder-подсказку 'Образец'."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            set_run_font(run, FONT_PRIMARY)
            run.font.size = TABLE_SLIDE_TITLE_FONT_SIZE
            run.font.color.rgb = COLOR_TEXT_PRIMARY
            return True
    return False


def _blank_other_placeholders(slide) -> None:
    """Опустошает прочие placeholder'ы, кроме title."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        try:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ""
        except Exception:
            pass


def _add_source_line(slide, report_date) -> None:
    """Подпись 'Источник: Cbonds, DD.MM.YYYY' справа снизу."""
    box = slide.shapes.add_textbox(
        Cm(22.00), TABLE_SLIDE_SOURCE_TOP, Cm(10.70), Cm(0.50),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Источник: Cbonds, {format_date_dmy(report_date)}"
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_SOURCE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_SECONDARY


def _add_footnote(slide) -> None:
    """Сноска об индикативности котировок."""
    box = slide.shapes.add_textbox(
        Cm(1.13), TABLE_SLIDE_FOOTNOTE_TOP, Cm(31.60), Cm(0.40),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Внимание! Приведенные котировки и расчеты являются индикативными "
        "и подлежат регулярному обновлению. Отдельные параметры могут "
        "существенно изменяться под влиянием рыночной конъюнктуры."
    )
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_FOOTNOTE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_MUTED


def _render_one_page(
    prs: PresentationType,
    page_bonds: List[Bond],
    coupon_formulas: Dict[str, str],
    offer_dates: Dict[str, str],
    highlighted: set,
    section_title: str,
    start_number: int,
    report_date,
) -> None:
    """Рендер одного слайда-страницы таблицы."""
    layout = prs.slide_layouts[_TABLE_SLIDE_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    _fill_title_placeholder(slide, SLIDE_TITLE)
    _blank_other_placeholders(slide)

    rows = [
        _bond_to_row(
            start_number + i, b,
            coupon_formulas.get(b.isin),
            offer_dates.get(b.isin),
            highlighted,
        )
        for i, b in enumerate(page_bonds)
    ]
    render_bond_table(
        slide,
        left=TABLE_SLIDE_LEFT,
        top=TABLE_SLIDE_TOP,
        columns=FLOATER_COLUMNS,
        rows=rows,
        section_title=section_title,
    )

    if report_date is not None:
        _add_source_line(slide, report_date)
    _add_footnote(slide)


# Максимум строк данных на один слайд.
# Подобрано эмпирически под размер слайда 33.87×19.05 см
# с учётом заголовка слайда, шапки секции+колонок, источника и сноски.
# Если попробовать больше — нижние строки наедут на дисклеймер.
FLOATERS_ROWS_PER_PAGE = 24


def render_floaters(
    prs: PresentationType,
    bonds: List[Bond],
    coupon_formulas: Dict[str, str],
    offer_dates: Optional[Dict[str, str]] = None,
    section_title: str = SECTION_TITLE_FLOATERS,
    report_date=None,
    rows_per_page: int = FLOATERS_ROWS_PER_PAGE,
    highlighted: Optional[set] = None,
) -> None:
    """
    Добавляет серию слайдов с таблицей флоатеров.

    Args:
        prs: презентация на основе template.pptx
        bonds: список облигаций
        coupon_formulas: словарь ISIN → формула купона (из bonds.yaml)
        offer_dates: словарь ISIN → дата оферты (DD.MM.YYYY)
        section_title: заголовок секции таблицы
        report_date: дата отчёта для подписи
        rows_per_page: строк данных на один слайд
        highlighted: множество ISIN, у которых в колонке "В перечне" ставится ✓
    """
    offer_dates = offer_dates or {}
    highlighted = highlighted or set()

    from datetime import date as _date
    sorted_bonds = sorted(
        bonds,
        key=lambda b: (b.matdate or _date.max),
    )

    total = len(sorted_bonds)
    if total == 0:
        return

    for page_idx in range(0, total, rows_per_page):
        page_bonds = sorted_bonds[page_idx : page_idx + rows_per_page]
        _render_one_page(
            prs,
            page_bonds=page_bonds,
            coupon_formulas=coupon_formulas,
            offer_dates=offer_dates,
            highlighted=highlighted,
            section_title=section_title,
            start_number=page_idx + 1,
            report_date=report_date,
        )
