"""
CLI-точка входа для сборки презентации.

Usage:
    python run.py                                 # дата = предыдущий торговый день
    python run.py --date 2026-04-16               # дата задана явно
    python run.py --date 2026-04-16 --source /tmp/bonds_mock.csv
"""
import argparse
import sys
from datetime import date, datetime, timezone, timedelta
from pathlib import Path

from pptx import Presentation

# добавляем src/ в путь, чтобы работали относительные импорты
sys.path.insert(0, str(Path(__file__).resolve().parent))

# Отчёт-специфичные настройки остаются в репозитории
from src.settings import (
    TEMPLATE_PATH,
    OUTPUT_DIR,
    CONFIG_DIR,
    REPORT_TITLE_LINES,
    COVER_DISCLAIMER,
)

# Общий код вынесен в пакет bonds_report_core
from bonds_report_core.sources import BONDS_DATA_URL
from bonds_report_core.config_loader import load_sections
from bonds_report_core.data_loader import load_bonds
from bonds_report_core.pdf_export import convert_pptx_to_pdf, LibreOfficeNotFoundError
from bonds_report_core.trading_day import previous_trading_day
from bonds_report_core.highlighted import load_highlighted_isins
from bonds_report_core.utils import remove_all_slides
from bonds_report_core.slides.cover import render_cover
from bonds_report_core.slides.glossary import render_glossary
from bonds_report_core.slides.disclaimer import render_disclaimer

# Табличные слайды специфичны для этого отчёта
from src.slides.bond_table_ofz import render_ofz_fixed
from src.slides.bond_table_floaters import render_floaters
from src.slides.bond_table_corp_fixed import render_corp_fixed


# Путь, куда будет сохраняться свежий PDF отчёта.
# Имя файла всегда одинаковое — PDF перезаписывается при каждом запуске.
# Так его удобно подхватывает presentation_bot на сервере.
PDF_OUTPUT_PATH = Path("/opt/presentation_bot/out/bonds_rub.pdf")

# Имя pptx-файла. Фиксированное — перезаписывается при каждом запуске,
# чтобы upstream-сервисы (пересылка, загрузка на диск и т.п.) могли
# ссылаться на один и тот же путь.
PPTX_FILENAME = "bonds_rub.pptx"


def parse_args():
    p = argparse.ArgumentParser(description="Сборка отчёта по облигациям в PowerPoint.")
    p.add_argument(
        "--date",
        type=lambda s: datetime.strptime(s, "%Y-%m-%d").date(),
        default=None,
        help=(
            "Дата отчёта в формате YYYY-MM-DD. По умолчанию — "
            "предыдущий торговый день (пн-пт)."
        ),
    )
    p.add_argument(
        "--source",
        default=BONDS_DATA_URL,
        help=(
            "Источник CSV с данными: URL или путь к локальному файлу. "
            f"По умолчанию — {BONDS_DATA_URL}"
        ),
    )
    p.add_argument(
        "--output",
        type=Path,
        default=None,
        help=(
            f"Путь к выходному .pptx. По умолчанию — output/{PPTX_FILENAME} "
            "(перезаписывается при каждом запуске)."
        ),
    )
    p.add_argument(
        "--pdf-output",
        type=Path,
        default=PDF_OUTPUT_PATH,
        help=(
            "Путь к выходному .pdf. По умолчанию — "
            f"{PDF_OUTPUT_PATH} (перезаписывается при каждом запуске)."
        ),
    )
    p.add_argument(
        "--skip-pdf",
        action="store_true",
        help="Не конвертировать в PDF (полезно для отладки на машине без LibreOffice).",
    )
    return p.parse_args()


def build_presentation(report_date, source: str, output_path: Path) -> Path:
    """Основной пайплайн сборки."""
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Не найден шаблон: {TEMPLATE_PATH}")

    date_only = report_date.date() if isinstance(report_date, datetime) else report_date

    # Читаем конфиг со списком ISIN.
    # Путь задаём явно (из settings), чтобы не зависеть от cwd.
    sections = load_sections(CONFIG_DIR / "bonds.yaml")
    sections_by_id = {s.id: s for s in sections}

    # Список ISIN, которые надо выделить галочкой "✓" в колонке "В перечне".
    # Файл /opt/config/highlighted_isins.csv — общий для обоих отчётов,
    # правится вручную. Если файла нет — highlighted=set() и галочек ни у кого.
    highlighted = load_highlighted_isins()
    if highlighted:
        print(f"Highlighted ISINs from перечня: {len(highlighted)}")

    # Открываем шаблон как основу. Удаляем 38 демо-слайдов — останется
    # только корпоративный стиль и layout'ы.
    prs = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(prs)

    # --- Слайд 1: обложка ---
    render_cover(prs, report_date, REPORT_TITLE_LINES, COVER_DISCLAIMER)

    # --- Слайд 2: ОФЗ с фиксированным купоном ---
    if "ofz_fixed" in sections_by_id:
        sec = sections_by_id["ofz_fixed"]
        print(f"Loading {len(sec.isins)} ОФЗ from {source}")
        bonds = load_bonds(sec.isins, source=source)
        render_ofz_fixed(prs, bonds, report_date=date_only, highlighted=highlighted)

    # --- Слайд 3: флоатеры ---
    if "floaters" in sections_by_id:
        sec = sections_by_id["floaters"]
        print(f"Loading {len(sec.isins)} флоатеров from {source}")
        bonds = load_bonds(sec.isins, source=source)
        render_floaters(
            prs,
            bonds,
            coupon_formulas=sec.coupon_formulas(),
            offer_dates=sec.offer_dates(),
            report_date=date_only,
            highlighted=highlighted,
        )

    # --- Слайд 4+: корпоративные фиксы ---
    if "corp_fixed" in sections_by_id:
        sec = sections_by_id["corp_fixed"]
        print(f"Loading {len(sec.isins)} корп-фиксов from {source}")
        bonds = load_bonds(sec.isins, source=source)
        render_corp_fixed(
            prs,
            bonds,
            offer_dates=sec.offer_dates(),
            report_date=date_only,
            highlighted=highlighted,
        )

    # --- Глоссарий (2 слайда) и дисклеймер (1 слайд) ---
    # Статический контент, не зависит от данных.
    render_glossary(prs)
    render_disclaimer(prs)

    # Здесь позже будут:
    # render_floaters(prs, ...)
    # render_corp_fixed(prs, ...)
    # render_glossary(prs)
    # render_disclaimer(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    args = parse_args()

    if args.date is None:
        args.date = datetime.now(tz=timezone(timedelta(hours=3)))

    if args.output is None:
        args.output = OUTPUT_DIR / PPTX_FILENAME

    print(f"Building report for {args.date} <- {args.source}")
    result = build_presentation(args.date, args.source, args.output)
    print(f"Saved: {result}")

    if args.skip_pdf:
        print("[pdf] Пропускаю конвертацию (--skip-pdf)")
        return

    # Экспорт в PDF. Ошибка конвертации не должна ронять всю сборку,
    # потому что pptx уже создан и это ценный результат — поэтому
    # ловим исключения и только пишем предупреждение.
    try:
        pdf_path = convert_pptx_to_pdf(args.output, args.pdf_output)
        print(f"PDF saved: {pdf_path}")
    except LibreOfficeNotFoundError as e:
        print(f"[pdf] Предупреждение: {e}")
    except Exception as e:
        print(f"[pdf] Ошибка конвертации: {e}")
        print(f"[pdf] pptx при этом сохранён: {args.output}")


if __name__ == "__main__":
    main()
